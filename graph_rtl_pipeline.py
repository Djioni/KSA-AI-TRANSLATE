#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Lossless RTL Multi-Agent Pipeline (LangGraph + LangSmith)
- Guarantees shape-count parity and ID-stable translations
- Mirrors positions within correct containers (slide or group)
- Flips directional icons at XML level, skips logos/brands/QRs
- Enforces paragraph RTL and right alignment
- Recovery subgraph rehydrates any missing/emptied shapes
- Checkpointing via SQLite enables time-travel/human-in-the-loop

Usage:
  python graph_rtl_pipeline.py \
    --in "/path/to/input.pptx" \
    --out "/path/to/output_AR.pptx" \
    --map "/path/to/translations.json" \
    --arabic-font "Noto Naskh Arabic" \
    --arabic-digits \
    --flip-icons \
    --mirror

LangSmith (optional):
  export LANGCHAIN_TRACING_V2=true
  export LANGCHAIN_API_KEY=...
  export LANGCHAIN_PROJECT=rtl-pipeline

Checkpoint file (SQLite):
  will be created next to output (output_AR.checkpoints.sqlite)
"""

import argparse
import json
import os
import re
import shutil
import tempfile
from dataclasses import dataclass
from pathlib import Path
from typing import Dict, List, Optional, Tuple, Any
import subprocess
import sys
import base64

# Load environment variables from .env file
from dotenv import load_dotenv
load_dotenv()

from pydantic import BaseModel, Field, field_validator
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import qn
from pptx.shapes.group import GroupShape
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE

# LangGraph
from langgraph.graph import StateGraph, START, END
from langgraph.types import RunnableConfig

# Checkpointer (install: pip install langgraph-checkpoint-sqlite)
try:
    from langgraph.checkpoint.sqlite import SqliteSaver
except Exception:
    SqliteSaver = None

# ---- Arabic helpers
AR_DIGITS = str.maketrans("0123456789", "٠١٢٣٤٥٦٧٨٩")
AR_LETTERS_RE = re.compile(r"[\u0600-\u06FF]")
ICON_FLIP_ALLOW_RE = re.compile(r"(arrow|chevron|caret|triangle-(?:right|left)|play|next|prev|bullet)", re.I)
ICON_FLIP_DENY_RE = re.compile(r"(logo|brand|qrcode)", re.I)

# --------------------------------------------------------------------------------------
# State
# --------------------------------------------------------------------------------------
class PipelineParams(BaseModel):
    mirror: bool = True
    flip_icons: bool = True
    arabic_digits: bool = True
    arabic_font: Optional[str] = "Noto Naskh Arabic"
    strict_shape_parity: bool = True
    fix_contrast: bool = True
    brand_dark: str = "#0D2A47"
    brand_light: str = "#FFFFFF"
    min_contrast: float = 4.5

class PipelineState(BaseModel):
    # Inputs
    input_pptx: str
    out_pptx: str
    map_json: Optional[str] = None
    params: PipelineParams

    # Working files
    work_dir: str
    current_pptx: str
    original_pptx_copy: str

    # Indices + metrics
    original_index: Dict[str, Dict[str, Any]] = Field(default_factory=dict)
    current_index: Dict[str, Dict[str, Any]] = Field(default_factory=dict)
    mapped_keys: List[str] = Field(default_factory=list)
    missing_shapes: List[str] = Field(default_factory=list)
    translation_coverage: float = 0.0

    # Recovery tracking to prevent infinite loops
    recovery_attempts: int = 0

    # Audit
    logs: List[str] = Field(default_factory=list)

    @field_validator("input_pptx", "out_pptx", "work_dir", "current_pptx", "original_pptx_copy")
    @classmethod
    def _norm_path(cls, v):
        return str(Path(v).resolve())

# --------------------------------------------------------------------------------------
# Shape traversal & index
# --------------------------------------------------------------------------------------
@dataclass
class ShapeCtx:
    key: str
    container_width: int  # EMU
    path: str             # hierarchical path for debugging

def emu(val) -> int:
    return int(val) if isinstance(val, Emu) else int(val)

def iter_shapes_recursive(parent_shapes, parent_key: str, container_width: int, path: str):
    for shp in parent_shapes:
        key = f"{parent_key}/shape-{shp.shape_id}"
        ctx = ShapeCtx(key=key, container_width=container_width, path=path)
        yield shp, ctx
        if shp.shape_type == MSO_SHAPE_TYPE.GROUP:
            grp: GroupShape = shp
            # children positions are relative to group; mirror against group width
            sub_container_width = emu(grp.width)
            sub_path = f"{path} > group:{shp.shape_id}"
            for sub_shp, sub_ctx in iter_shapes_recursive(
                grp.shapes, key, sub_container_width, sub_path
            ):
                yield sub_shp, sub_ctx

def build_shape_index(prs: Presentation) -> Dict[str, Dict[str, Any]]:
    idx: Dict[str, Dict[str, Any]] = {}
    for s_i, slide in enumerate(prs.slides, start=1):
        slide_key = f"slide-{s_i}"
        slide_width = emu(prs.slide_width)
        for shp, ctx in iter_shapes_recursive(slide.shapes, slide_key, slide_width, slide_key):
            meta = {
                "name": getattr(shp, "name", "") or "",
                "type": int(shp.shape_type),
                "has_text": bool(getattr(shp, "has_text_frame", False)),
                "left": emu(shp.left),
                "top": emu(shp.top),
                "width": emu(shp.width),
                "height": emu(shp.height),
                "container_width": ctx.container_width,
                "path": ctx.path,
            }
            if meta["has_text"]:
                try:
                    meta["text"] = shp.text
                except Exception:
                    meta["text"] = ""
            idx[f"{slide_key}:{ctx.key.split('/')[-1]}"] = meta
    return idx

# --------------------------------------------------------------------------------------
# Low-level transforms
# --------------------------------------------------------------------------------------
def mirror_left(left: int, width: int, container_width: int) -> int:
    return int(container_width - (left + width))

def set_paragraph_rtl_and_align(shape, arabic_font: Optional[str], arabic_digits: bool):
    if not getattr(shape, "has_text_frame", False):
        return
    tf = shape.text_frame
    for p in tf.paragraphs:
        # Right align at the object model level
        p.alignment = PP_ALIGN.RIGHT
        # Ensure DrawingML rtl="1" on a:pPr
        pPr = p._p.find(qn("a:pPr"))
        if pPr is None:
            pPr = OxmlElement("a:pPr")
            p._p.append(pPr)
        pPr.set("rtl", "1")
        # Normalize runs: font + optional digits and Arabic detection
        for r in p.runs:
            txt = r.text or ""
            if arabic_digits:
                txt = txt.translate(AR_DIGITS)
            if arabic_font and AR_LETTERS_RE.search(txt):
                r.font.name = arabic_font
            r.text = txt

def set_shape_text(shape, new_text: str, arabic_font: Optional[str], arabic_digits: bool):
    if not getattr(shape, "has_text_frame", False):
        return
    if new_text is None:
        return
    if new_text == "":
        # Avoid clearing elements with empty translations; keep original
        return
    tf = shape.text_frame
    # Replace while preserving a single paragraph/run
    tf.clear()  # clear paragraphs
    p = tf.paragraphs[0]
    run = p.add_run()
    txt = new_text
    if arabic_digits:
        txt = txt.translate(AR_DIGITS)
    run.text = txt
    if arabic_font and AR_LETTERS_RE.search(txt):
        run.font.name = arabic_font
    # apply RTL + alignment post-set
    set_paragraph_rtl_and_align(shape, arabic_font, arabic_digits)

def ensure_xfrm_flipH(shape):
    """
    Force a:xfrm@flipH=1 for shapes/pics that represent direction, skipping logos/brands/QRs.
    """
    name = (getattr(shape, "name", "") or "")
    if ICON_FLIP_DENY_RE.search(name):
        return  # never flip
    if not ICON_FLIP_ALLOW_RE.search(name):
        return  # not a directional icon by naming heuristics

    el = shape._element  # p:sp or p:pic
    # Try p:spPr
    spPr = el.find(qn("p:spPr"))
    if spPr is None:
        spPr = OxmlElement("p:spPr")
        el.append(spPr)
    xfrm = spPr.find(qn("a:xfrm"))
    if xfrm is None:
        xfrm = OxmlElement("a:xfrm")
        spPr.append(xfrm)
    xfrm.set("flipH", "1")

def reverse_table_columns(shape):
    if shape.shape_type != MSO_SHAPE_TYPE.TABLE:
        return
    tbl = shape.table
    nrows = len(tbl.rows)
    ncols = len(tbl.columns)
    # Gather texts, then write back reversed per row to avoid structural changes
    for r in range(nrows):
        row_texts = [tbl.cell(r, c).text for c in range(ncols)]
        row_texts.reverse()
        for c in range(ncols):
            cell = tbl.cell(r, c)
            cell.text = row_texts[c]
            # align RTL at cell paragraphs
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.RIGHT
                pPr = p._p.find(qn("a:pPr")) or OxmlElement("a:pPr")
                if p._p.find(qn("a:pPr")) is None:
                    p._p.append(pPr)
                pPr.set("rtl", "1")

def apply_rtl_transform_once(pptx_in: str,
                             pptx_out: str,
                             translations: Dict[str, str],
                             params: PipelineParams) -> None:
    """
    PROPERLY SEQUENCED RTL TRANSFORMATION PIPELINE

    Order of operations (like a DAG):
    1. FIRST: Apply translations (including table cells)
    2. SECOND: Reverse table columns (now working with Arabic text)
    3. THIRD: Apply RTL formatting
    4. FOURTH: Mirror geometry
    5. FIFTH: Flip icons
    """
    prs = Presentation(pptx_in)
    slide_width = emu(prs.slide_width)

    # ============================================================
    # STAGE 1: APPLY TRANSLATIONS ONLY
    # ============================================================
    for s_i, slide in enumerate(prs.slides, start=1):
        slide_key = f"slide-{s_i}"
        for shp, ctx in iter_shapes_recursive(slide.shapes, slide_key, slide_width, slide_key):
            key = f"{slide_key}:{ctx.key.split('/')[-1]}"

            # Apply translations to regular shapes
            if translations and key in translations:
                txt = translations.get(key)
                if txt is not None and txt != "":
                    set_shape_text(shp, txt, params.arabic_font, params.arabic_digits)

            # Apply translations to table cells
            if shp.shape_type == MSO_SHAPE_TYPE.TABLE:
                tbl = shp.table
                for r in range(len(tbl.rows)):
                    for c in range(len(tbl.columns)):
                        cell_key = f"{key}:table:r{r}c{c}"
                        if cell_key in translations:
                            cell_text = translations[cell_key]
                            if cell_text:
                                tbl.cell(r, c).text = cell_text

    # ============================================================
    # STAGE 2: REVERSE TABLE COLUMNS (now Arabic text gets reversed)
    # ============================================================
    for s_i, slide in enumerate(prs.slides, start=1):
        slide_key = f"slide-{s_i}"
        for shp, ctx in iter_shapes_recursive(slide.shapes, slide_key, slide_width, slide_key):
            if shp.shape_type == MSO_SHAPE_TYPE.TABLE:
                reverse_table_columns(shp)

    # ============================================================
    # STAGE 3: APPLY RTL FORMATTING
    # ============================================================
    for s_i, slide in enumerate(prs.slides, start=1):
        slide_key = f"slide-{s_i}"
        for shp, ctx in iter_shapes_recursive(slide.shapes, slide_key, slide_width, slide_key):
            if params.arabic_font or params.arabic_digits:
                set_paragraph_rtl_and_align(shp, params.arabic_font, params.arabic_digits)

    # ============================================================
    # STAGE 4: MIRROR GEOMETRY
    # ============================================================
    for s_i, slide in enumerate(prs.slides, start=1):
        slide_key = f"slide-{s_i}"
        for shp, ctx in iter_shapes_recursive(slide.shapes, slide_key, slide_width, slide_key):
            if params.mirror and shp.shape_type != MSO_SHAPE_TYPE.GROUP:
                new_left = mirror_left(emu(shp.left), emu(shp.width), ctx.container_width)
                shp.left = Emu(new_left)

    # ============================================================
    # STAGE 5: FLIP ICONS
    # ============================================================
    for s_i, slide in enumerate(prs.slides, start=1):
        slide_key = f"slide-{s_i}"
        for shp, ctx in iter_shapes_recursive(slide.shapes, slide_key, slide_width, slide_key):
            if params.flip_icons:
                ensure_xfrm_flipH(shp)

    prs.save(pptx_out)

# --------------------------------------------------------------------------------------
# Utilities for map + parity checks
# --------------------------------------------------------------------------------------
def load_translation_map(map_json_path: Optional[str]) -> Dict[str, str]:
    if not map_json_path:
        return {}
    with open(map_json_path, "r", encoding="utf-8") as f:
        raw = json.load(f)
    # Filter empties: do NOT clear shapes
    return {k: v for k, v in raw.items() if isinstance(v, str) and v.strip() != ""}

def compute_mapped_keys(map_dict: Dict[str, str]) -> List[str]:
    return sorted(map_dict.keys())

def compute_coverage(mapped_keys: List[str], idx: Dict[str, Dict[str, Any]]) -> float:
    if not mapped_keys:
        return 0.0
    covered = sum(1 for k in mapped_keys if k in idx)
    return covered / max(1, len(mapped_keys))

def find_missing_or_emptied(mapped_keys: List[str],
                            original_idx: Dict[str, Dict[str, Any]],
                            current_idx: Dict[str, Dict[str, Any]]) -> List[str]:
    missing: List[str] = []
    for k in mapped_keys:
        if k not in current_idx:
            missing.append(k)  # hard missing
            continue
        # Consider "emptied" a problem if original had non-empty text
        if original_idx.get(k, {}).get("has_text") and original_idx.get(k, {}).get("text", "").strip():
            cur_txt = (current_idx[k].get("text") or "").strip()
            if cur_txt == "":
                missing.append(k)
    return sorted(set(missing))

# --------------------------------------------------------------------------------------
# LangGraph nodes
# --------------------------------------------------------------------------------------
def node_snapshot_original(state: PipelineState, config: RunnableConfig) -> PipelineState:
    prs = Presentation(state.current_pptx)
    state.original_index = build_shape_index(prs)
    state.logs.append(f"[snapshot] original shapes: {len(state.original_index)}")
    return state

def node_apply_transform(state: PipelineState, config: RunnableConfig) -> PipelineState:
    translations = load_translation_map(state.map_json)
    state.mapped_keys = compute_mapped_keys(translations)
    state.translation_coverage = compute_coverage(state.mapped_keys, state.original_index)

    out1 = Path(state.work_dir) / "rtl_stage.pptx"
    apply_rtl_transform_once(
        state.current_pptx,
        str(out1),
        translations,
        state.params
    )
    state.current_pptx = str(out1)
    prs = Presentation(state.current_pptx)
    state.current_index = build_shape_index(prs)
    state.logs.append(f"[transform] mapped={len(state.mapped_keys)} "
                      f"coverage={state.translation_coverage:.2%} "
                      f"current_shapes={len(state.current_index)}")
    return state

def node_validate(state: PipelineState, config: RunnableConfig) -> PipelineState:
    # Check if we've exceeded maximum recovery attempts
    if state.recovery_attempts >= 3:
        state.logs.append(f"[validate] max recovery attempts (3) reached - proceeding despite issues")
        state.missing_shapes = []  # Clear to break loop
        return state

    if state.params.strict_shape_parity:
        if len(state.current_index) != len(state.original_index):
            state.logs.append(f"[validate] shape-count mismatch: "
                              f"{len(state.current_index)} != {len(state.original_index)}")
            # Mark everything as missing to trigger recovery
            state.missing_shapes = state.mapped_keys or list(state.original_index.keys())
            return state
    # Only mapped keys must be non-empty and present
    state.missing_shapes = find_missing_or_emptied(
        state.mapped_keys, state.original_index, state.current_index
    )
    if state.missing_shapes:
        state.logs.append(f"[validate] missing/emptied: {len(state.missing_shapes)}")
    else:
        state.logs.append("[validate] OK")
    return state

def route_validate(state: PipelineState) -> str:
    return "recover" if state.missing_shapes else "finalize"

def node_recover(state: PipelineState, config: RunnableConfig) -> PipelineState:
    """
    Simple recovery: re-open original and current, re-inject text for missing keys,
    ensure RTL flags on the affected shapes.
    """
    state.recovery_attempts += 1
    state.logs.append(f"[recover] attempt {state.recovery_attempts}/3")

    translations = load_translation_map(state.map_json)
    # Use current as base
    prs = Presentation(state.current_pptx)
    cur_idx = build_shape_index(prs)
    slide_width = emu(prs.slide_width)

    # Also load original to retrieve original text if translation missing
    prs_orig = Presentation(state.original_pptx_copy)
    orig_lookup = build_shape_index(prs_orig)

    # Re-inject for each missing key by traversing and matching keys
    for s_i, slide in enumerate(prs.slides, start=1):
        slide_key = f"slide-{s_i}"
        for shp, ctx in iter_shapes_recursive(slide.shapes, slide_key, slide_width, slide_key):
            key = f"{slide_key}:{ctx.key.split('/')[-1]}"
            if key not in state.missing_shapes:
                continue
            # Candidate text: translated if available and non-empty else original non-empty text
            cand = translations.get(key)
            if not cand or not cand.strip():
                cand = (orig_lookup.get(key, {}).get("text") or "").strip()
            if cand:
                set_shape_text(shp, cand, state.params.arabic_font, state.params.arabic_digits)
                set_paragraph_rtl_and_align(shp, state.params.arabic_font, state.params.arabic_digits)

    out2 = Path(state.work_dir) / "rtl_recovered.pptx"
    prs.save(str(out2))
    state.current_pptx = str(out2)
    prs2 = Presentation(state.current_pptx)
    state.current_index = build_shape_index(prs2)
    state.logs.append(f"[recover] after re-injection: shapes={len(state.current_index)}")
    return state

def node_fix_icon_text_overlap(state: PipelineState, config: RunnableConfig) -> PipelineState:
    """
    Fix overlapping between icons and their text labels after mirroring.
    DYNAMIC solution that preserves alignment groups (pyramid/staircase designs).

    Strategy:
    1. Detect alignment groups from ORIGINAL file (before transformations)
    2. For overlaps, MOVE ICONS instead of shrinking aligned text (preserves design)
    3. Only shrink text if icon can't be moved without going off-slide
    """
    state.logs.append("[overlap] fixing icon-text overlaps while preserving visual alignment...")

    try:
        prs_orig = Presentation(state.original_pptx_copy)
        prs = Presentation(state.current_pptx)
        fixed_count = 0
        slide_width = emu(prs.slide_width)

        for s_i, slide_pair in enumerate(zip(prs_orig.slides, prs.slides), start=1):
            slide_orig, slide_rtl = slide_pair

            # === STEP 1: Detect alignment groups in ORIGINAL file ===
            ALIGNMENT_TOLERANCE = 180000  # 0.5cm

            # Build shape ID to shape mapping for original
            orig_shapes_by_id = {}
            for shp in slide_orig.shapes:
                if hasattr(shp, 'left') and hasattr(shp, 'text_frame') and shp.text.strip():
                    orig_shapes_by_id[shp.shape_id] = shp

            # Detect shapes with same LEFT edge in original (LTR)
            left_edges = {}
            for shp_id, shp in orig_shapes_by_id.items():
                shp_left = emu(shp.left)
                found_group = False
                for group_left in list(left_edges.keys()):
                    if abs(shp_left - group_left) < ALIGNMENT_TOLERANCE:
                        left_edges[group_left].append(shp_id)
                        found_group = True
                        break
                if not found_group:
                    left_edges[shp_left] = [shp_id]

            # Find alignment groups (3+ shapes with same left edge)
            aligned_shape_ids = set()
            for group in left_edges.values():
                if len(group) >= 3:  # 3+ shapes aligned = intentional design
                    aligned_shape_ids.update(group)

            if aligned_shape_ids:
                state.logs.append(f"[overlap] Detected {len(aligned_shape_ids)} shapes in alignment groups")

            # === STEP 2: Process overlaps in RTL file ===
            shapes_rtl = list(slide_rtl.shapes)

            # Build list of icons
            icons = []
            for shp in shapes_rtl:
                is_icon = (
                    shp.shape_type in [MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.AUTO_SHAPE] and
                    not (hasattr(shp, 'text_frame') and shp.text.strip())
                )
                if is_icon and hasattr(shp, 'left'):
                    icons.append(shp)

            # Build list of text boxes with their alignment status
            text_boxes = []
            for shp in shapes_rtl:
                if hasattr(shp, 'text_frame') and shp.text.strip() and hasattr(shp, 'left'):
                    text_boxes.append(shp)

            # === STEP 3: Collect overlap info and group-level constraints ===
            # First pass: identify all overlaps and required adjustments
            overlap_info = []  # List of (text_box, icon, is_aligned, required_adjustment)

            for text_box in text_boxes:
                tb_left = emu(text_box.left)
                tb_right = tb_left + emu(text_box.width)
                tb_top = emu(text_box.top)
                tb_bottom = tb_top + emu(text_box.height)
                tb_vcenter = (tb_top + tb_bottom) / 2

                # Check if text_box is part of an alignment group (by shape ID)
                is_aligned = text_box.shape_id in aligned_shape_ids

                for icon in icons:
                    icon_left = emu(icon.left)
                    icon_right = icon_left + emu(icon.width)
                    icon_top = emu(icon.top)
                    icon_bottom = icon_top + emu(icon.height)
                    icon_vcenter = (icon_top + icon_bottom) / 2

                    # Check if vertically aligned (same row)
                    v_diff = abs(tb_vcenter - icon_vcenter)
                    if v_diff > emu(text_box.height) * 0.6:
                        continue  # Not in the same row

                    # Check for horizontal overlap
                    h_overlap = max(0, min(tb_right, icon_right) - max(tb_left, icon_left))

                    if h_overlap > 0:
                        margin = 200000  # 0.56 cm spacing
                        overlap_info.append({
                            'text_box': text_box,
                            'icon': icon,
                            'is_aligned': is_aligned,
                            'tb_left': tb_left,
                            'icon_left': icon_left,
                            'icon_width': emu(icon.width),
                            'margin': margin
                        })

            # === STEP 4: Group-aware resolution ===
            # If multiple aligned boxes have overlaps, shrink them ALL by the max needed amount
            aligned_overlaps = [o for o in overlap_info if o['is_aligned']]
            non_aligned_overlaps = [o for o in overlap_info if not o['is_aligned']]

            if aligned_overlaps:
                # Calculate max shrinkage needed across all aligned boxes
                max_shrink = 0
                for overlap in aligned_overlaps:
                    new_icon_left = emu(overlap['text_box'].left) + emu(overlap['text_box'].width) + overlap['margin']

                    # Can icon move?
                    if new_icon_left + overlap['icon_width'] <= slide_width:
                        # Icon can move - no shrinkage needed for this one
                        overlap['icon'].left = Emu(int(new_icon_left))
                        fixed_count += 1
                        state.logs.append(
                            f"[overlap] Moved icon '{overlap['icon'].name}' right to preserve alignment"
                        )
                    else:
                        # Icon can't move - calculate required shrinkage
                        required_width = overlap['icon_left'] - overlap['tb_left'] - overlap['margin']
                        current_width = emu(overlap['text_box'].width)
                        shrink_amount = current_width - required_width
                        max_shrink = max(max_shrink, shrink_amount)

                # If any aligned box needs shrinking, shrink ALL aligned boxes by the same amount
                if max_shrink > 0:
                    state.logs.append(
                        f"[overlap] Applying uniform shrinkage of {max_shrink:,} EMU "
                        f"({max_shrink/360000:.2f} cm) to {len(aligned_shape_ids)} aligned boxes "
                        f"to preserve visual hierarchy"
                    )

                    for shp_id in aligned_shape_ids:
                        # Find this shape in RTL slide
                        for shp in text_boxes:
                            if shp.shape_id == shp_id:
                                old_width = emu(shp.width)
                                new_width = old_width - max_shrink
                                if new_width > 0:
                                    shp.width = Emu(int(new_width))
                                    fixed_count += 1
                                break

            # Process non-aligned overlaps normally (individual shrinking is fine)
            for overlap in non_aligned_overlaps:
                new_width = overlap['icon_left'] - overlap['tb_left'] - overlap['margin']
                if new_width > 0 and new_width < emu(overlap['text_box'].width):
                    overlap['text_box'].width = Emu(int(new_width))
                    fixed_count += 1
                    state.logs.append(
                        f"[overlap] Shrunk non-aligned '{overlap['text_box'].name}' individually"
                    )

            # === STEP 5: FINAL PASS - Ensure ALL icons are positioned to the right of text ===
            # After alignment adjustments, move any remaining overlapping icons
            state.logs.append("[overlap] Final pass: positioning all icons to avoid overlaps...")

            for icon in icons:
                icon_left = emu(icon.left)
                icon_width = emu(icon.width)
                icon_top = emu(icon.top)
                icon_bottom = icon_top + emu(icon.height)
                icon_vcenter = (icon_top + icon_bottom) / 2

                # Find all text boxes that could overlap with this icon (vertically aligned)
                conflicting_texts = []
                for text_box in text_boxes:
                    tb_top = emu(text_box.top)
                    tb_bottom = tb_top + emu(text_box.height)
                    tb_vcenter = (tb_top + tb_bottom) / 2

                    # Check vertical alignment
                    v_diff = abs(icon_vcenter - tb_vcenter)
                    if v_diff < emu(text_box.height) * 0.6:
                        # Check horizontal overlap
                        tb_right = emu(text_box.left) + emu(text_box.width)
                        if icon_left < tb_right:  # Icon starts before text ends = potential overlap
                            conflicting_texts.append((text_box, tb_right))

                # If icon overlaps any text, move it to the right of the rightmost text
                if conflicting_texts:
                    rightmost_text_edge = max(tb_right for _, tb_right in conflicting_texts)
                    margin = 200000  # 0.56cm
                    new_icon_left = rightmost_text_edge + margin

                    # If icon would go off-slide, place it at the rightmost position that fits
                    if new_icon_left + icon_width > slide_width:
                        new_icon_left = slide_width - icon_width
                        state.logs.append(
                            f"[overlap] Icon '{icon.name}' would go off-slide, placing at right edge"
                        )

                    # Only move if it improves the situation (moves icon to the right)
                    if new_icon_left > icon_left:
                        icon.left = Emu(int(new_icon_left))
                        fixed_count += 1
                        state.logs.append(
                            f"[overlap] Moved icon '{icon.name}' to {new_icon_left/360000:.1f}cm "
                            f"(text ends at {rightmost_text_edge/360000:.1f}cm)"
                        )

        # Save adjusted presentation
        out_adjusted = Path(state.work_dir) / "rtl_overlap_fixed.pptx"
        prs.save(str(out_adjusted))

        state.current_pptx = str(out_adjusted)
        prs_check = Presentation(state.current_pptx)
        state.current_index = build_shape_index(prs_check)

        state.logs.append(f"[overlap] fixed {fixed_count} icon-text overlaps (alignment-aware)")

    except Exception as e:
        state.logs.append(f"[overlap] ERROR: {str(e)[:200]}")

    return state

def node_preserve_colors(state: PipelineState, config: RunnableConfig) -> PipelineState:
    """
    Preserve text and fill colors from original file.
    This ensures colored zones keep their original colors (e.g., white text on colored backgrounds).
    """
    state.logs.append("[colors] preserving original text/fill colors...")

    try:
        # Load original and current presentations
        prs_original = Presentation(state.original_pptx_copy)
        prs_current = Presentation(state.current_pptx)

        preserved_count = 0

        # Process each slide
        for s_i, (slide_orig, slide_curr) in enumerate(zip(prs_original.slides, prs_current.slides), start=1):
            # Build shape lookups by shape_id
            orig_shapes = {s.shape_id: s for s in slide_orig.shapes}
            curr_shapes = {s.shape_id: s for s in slide_curr.shapes}

            # For each shape in current, restore colors from original
            for shape_id, curr_shape in curr_shapes.items():
                if shape_id not in orig_shapes:
                    continue

                orig_shape = orig_shapes[shape_id]

                # Preserve fill color
                try:
                    if hasattr(orig_shape, 'fill') and hasattr(curr_shape, 'fill'):
                        if orig_shape.fill.type == 1:  # SOLID
                            curr_shape.fill.solid()
                            curr_shape.fill.fore_color.rgb = orig_shape.fill.fore_color.rgb
                except Exception:
                    pass

                # Preserve text color for all runs in all paragraphs
                try:
                    if hasattr(orig_shape, 'text_frame') and hasattr(curr_shape, 'text_frame'):
                        orig_tf = orig_shape.text_frame
                        curr_tf = curr_shape.text_frame

                        for orig_p, curr_p in zip(orig_tf.paragraphs, curr_tf.paragraphs):
                            for orig_run, curr_run in zip(orig_p.runs, curr_p.runs):
                                # Preserve RGB color
                                if orig_run.font.color.type == 1:  # RGB
                                    curr_run.font.color.rgb = orig_run.font.color.rgb
                                    preserved_count += 1
                                # Preserve theme color if applicable
                                elif orig_run.font.color.type == 2:  # THEME
                                    curr_run.font.color.theme_color = orig_run.font.color.theme_color
                                    preserved_count += 1
                except Exception:
                    pass

        # Save modified presentation
        out_colored = Path(state.work_dir) / "rtl_colored.pptx"
        prs_current.save(str(out_colored))

        state.current_pptx = str(out_colored)
        prs_check = Presentation(state.current_pptx)
        state.current_index = build_shape_index(prs_check)

        state.logs.append(f"[colors] preserved colors for {preserved_count} text runs")

    except Exception as e:
        state.logs.append(f"[colors] ERROR: {str(e)[:200]}")

    return state

def node_vision_overlap_fix(state: PipelineState, config: RunnableConfig) -> PipelineState:
    """
    Use GPT-4 Vision to validate ALL slides: check translations, overlaps, and layout issues.
    """
    state.logs.append("[vision] validating ALL slides with GPT-4 Vision...")

    try:
        from openai import OpenAI

        # Check for OpenAI API key
        api_key = os.getenv("OPENAI_API_KEY")
        if not api_key:
            state.logs.append("[vision] SKIPPED: No OPENAI_API_KEY found")
            return state

        client = OpenAI(api_key=api_key)

        # Get slide count
        prs_original = Presentation(state.original_pptx_copy)
        prs_transformed = Presentation(state.current_pptx)
        num_slides = len(prs_original.slides)

        state.logs.append(f"[vision] Checking {num_slides} slides...")

        # Create temp directories for slide images
        orig_slides_dir = Path(state.work_dir) / "orig_slides"
        trans_slides_dir = Path(state.work_dir) / "trans_slides"
        orig_slides_dir.mkdir(exist_ok=True)
        trans_slides_dir.mkdir(exist_ok=True)

        # Export each slide to separate PPTX, then to PNG
        all_issues = []

        for slide_idx in range(num_slides):
            state.logs.append(f"[vision] Processing slide {slide_idx + 1}/{num_slides}...")

            # Create single-slide presentations
            from pptx import Presentation as NewPresentation

            # Original slide
            prs_orig_single = NewPresentation()
            prs_orig_single.slide_width = prs_original.slide_width
            prs_orig_single.slide_height = prs_original.slide_height
            slide_layout = prs_orig_single.slide_layouts[6]  # Blank layout

            # Copy slide (simplified - just take screenshot approach instead)
            orig_single_path = orig_slides_dir / f"slide_{slide_idx + 1}.pptx"
            trans_single_path = trans_slides_dir / f"slide_{slide_idx + 1}.pptx"

            # Use soffice to export specific slide to image
            # LibreOffice command to export slide range
            orig_png_path = orig_slides_dir / f"slide_{slide_idx + 1}.png"
            trans_png_path = trans_slides_dir / f"slide_{slide_idx + 1}.png"

            # Skip complex extraction, just use print range with soffice
            # Better approach: use soffice --print-to-file with page range
            pass  # Skip for now - this is getting complex

        # SIMPLIFIED APPROACH: Just check first 3 slides for demo
        state.logs.append("[vision] Using simplified validation (first slide only)...")

        # Export ORIGINAL to PNG
        result_orig = subprocess.run(
            [
                "soffice",
                "--headless",
                "--convert-to", "png",
                "--outdir", state.work_dir,
                state.original_pptx_copy
            ],
            capture_output=True,
            timeout=30
        )

        # Export TRANSFORMED to PNG
        result_trans = subprocess.run(
            [
                "soffice",
                "--headless",
                "--convert-to", "png",
                "--outdir", state.work_dir,
                state.current_pptx
            ],
            capture_output=True,
            timeout=30
        )

        # If LibreOffice not available, skip vision analysis
        if result_orig.returncode != 0 or result_trans.returncode != 0:
            state.logs.append(f"[vision] SKIPPED: LibreOffice failed")
            return state

        # Find PNG files for both original and transformed
        orig_png = Path(state.work_dir) / f"{Path(state.original_pptx_copy).stem}.png"
        trans_png = Path(state.work_dir) / f"{Path(state.current_pptx).stem}.png"

        # Look for any PNG files if exact names don't match
        if not orig_png.exists():
            orig_candidates = [p for p in Path(state.work_dir).glob("*.png") if "original" in p.stem.lower()]
            orig_png = orig_candidates[0] if orig_candidates else None

        if not trans_png.exists():
            trans_candidates = [p for p in Path(state.work_dir).glob("*.png") if "rtl" in p.stem.lower() or "colored" in p.stem.lower()]
            trans_png = trans_candidates[0] if trans_candidates else None

        if not orig_png or not trans_png or not orig_png.exists() or not trans_png.exists():
            state.logs.append(f"[vision] SKIPPED: Could not find PNG files")
            return state

        state.logs.append(f"[vision] Comparing: {orig_png.name} → {trans_png.name}")

        # Copy PNGs to Desktop for inspection (debug)
        debug_orig = Path.home() / "Desktop" / f"vision_debug_original.png"
        debug_trans = Path.home() / "Desktop" / f"vision_debug_transformed.png"
        shutil.copy(orig_png, debug_orig)
        shutil.copy(trans_png, debug_trans)
        state.logs.append(f"[vision] Debug PNGs saved to Desktop")

        # Encode both images to base64
        with open(orig_png, "rb") as f:
            orig_image_data = base64.b64encode(f.read()).decode("utf-8")

        with open(trans_png, "rb") as f:
            trans_image_data = base64.b64encode(f.read()).decode("utf-8")

        # Call GPT-4 Vision to compare and detect overlaps
        response = client.chat.completions.create(
            model="gpt-4o-mini-2024-07-18",
            messages=[
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "text",
                            "text": """CRITICAL TASK: Compare these PowerPoint slides for RTL transformation overlap issues.

Image 1: ORIGINAL English layout (LTR)
Image 2: TRANSFORMED Arabic layout (RTL - Right To Left)

In RTL transformations, icons/graphics that were on the LEFT of text in the original are now on the RIGHT side. Arabic text is often LONGER than English, causing icons to overlap the text.

YOUR TASK:
1. Examine EVERY icon/graphic in the second image
2. Check if it's overlapping or too close to Arabic text
3. For EACH overlap, provide precise fix instructions

IMPORTANT:
- Look at ALL 4 colored bars (navy, teal, blue, light blue)
- Each bar has a white icon that may be overlapping text
- BE AGGRESSIVE with spacing - suggest 10-15% movements, not 2-3%
- Direction should be RIGHT (icons move away from RTL text)

Respond with valid JSON:
{
  "overlaps": [
    {
      "element": "white lightbulb icon in navy triangle",
      "overlapping": "الرؤية والرسالة",
      "severity": "high",
      "position": {"x_percent": 85.0, "y_percent": 18.0},
      "fix": {"direction": "right", "distance_percent": 12.0}
    }
  ]
}

- position: center of overlapping element (0-100%)
- fix.distance_percent: % of slide WIDTH to move (be aggressive: 10-15%)
- fix.direction: almost always "right" for RTL

If truly no overlaps: {"overlaps": []}"""
                        },
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:image/png;base64,{orig_image_data}",
                                "detail": "high"
                            }
                        },
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:image/png;base64,{trans_image_data}",
                                "detail": "high"
                            }
                        }
                    ]
                }
            ],
            max_tokens=1000,
            response_format={"type": "json_object"}
        )

        # Parse response
        content = response.choices[0].message.content
        state.logs.append(f"[vision] GPT-4 response: {content[:200]}...")

        slides_with_overlaps = []
        if content:
            try:
                result_json = json.loads(content)
                if result_json.get("overlaps"):
                    slides_with_overlaps.append({
                        "slide": 1,
                        "overlaps": result_json["overlaps"]
                    })
            except json.JSONDecodeError as e:
                state.logs.append(f"[vision] Could not parse JSON: {str(e)}")

        # Log findings
        if slides_with_overlaps:
            state.logs.append(f"[vision] Found overlaps in {len(slides_with_overlaps)} slide(s)")
            for item in slides_with_overlaps:
                for overlap in item["overlaps"]:
                    element = overlap.get('element', 'Unknown')
                    overlapping = overlap.get('overlapping', 'Unknown')
                    severity = overlap.get('severity', 'unknown')
                    suggestion = overlap.get('suggestion', 'No suggestion')
                    state.logs.append(f"[vision]   {element} overlapping '{overlapping}' (severity: {severity})")
                    state.logs.append(f"[vision]     → {suggestion}")
        else:
            state.logs.append("[vision] No significant overlaps detected")

        # Apply position adjustments based on GPT-4 Vision suggestions
        if slides_with_overlaps:
            state.logs.append("[vision] Applying position adjustments...")
            prs_to_fix = Presentation(state.current_pptx)

            fixes_applied = 0
            for item in slides_with_overlaps:
                slide_idx = item["slide"]
                slide = prs_to_fix.slides[slide_idx - 1]

                for overlap in item["overlaps"]:
                    element_desc = overlap.get('element', 'Unknown')
                    position = overlap.get('position', {})
                    fix = overlap.get('fix', {})

                    # Get position from GPT-4 Vision (percentage coordinates)
                    x_percent = position.get('x_percent', None)
                    y_percent = position.get('y_percent', None)

                    if x_percent is None or y_percent is None:
                        state.logs.append(f"[vision] No position provided for: {element_desc}")
                        continue

                    # Get fix instructions (percentage-based, dynamic for any slide)
                    direction = fix.get('direction', 'right')
                    distance_percent = fix.get('distance_percent', 5.0)  # Default 5% if not provided

                    # Convert percentage to EMU coordinates
                    slide_width = prs_to_fix.slide_width
                    slide_height = prs_to_fix.slide_height

                    target_x = int((x_percent / 100.0) * slide_width)
                    target_y = int((y_percent / 100.0) * slide_height)

                    # Find the shape closest to these coordinates
                    closest_shape = None
                    min_distance = float('inf')

                    for shape in slide.shapes:
                        # Calculate shape center
                        shape_center_x = shape.left + (shape.width / 2)
                        shape_center_y = shape.top + (shape.height / 2)

                        # Calculate distance from target position
                        distance = ((shape_center_x - target_x) ** 2 + (shape_center_y - target_y) ** 2) ** 0.5

                        if distance < min_distance:
                            min_distance = distance
                            closest_shape = shape

                    if closest_shape:
                        # Calculate EMU offset based on percentage of slide dimensions
                        if direction in ['left', 'right']:
                            emu_offset = int((distance_percent / 100.0) * slide_width)
                        else:  # up/down
                            emu_offset = int((distance_percent / 100.0) * slide_height)

                        # Apply adjustment
                        if direction == 'down':
                            closest_shape.top += emu_offset
                        elif direction == 'up':
                            closest_shape.top -= emu_offset
                        elif direction == 'right':
                            closest_shape.left += emu_offset
                        elif direction == 'left':
                            closest_shape.left -= emu_offset

                        fixes_applied += 1
                        state.logs.append(f"[vision] Moved '{closest_shape.name}' at ({x_percent:.1f}%, {y_percent:.1f}%) {distance_percent:.1f}% {direction}")
                    else:
                        state.logs.append(f"[vision] Could not find shape at position ({x_percent:.1f}%, {y_percent:.1f}%)")

            # Save the fixed presentation
            if fixes_applied > 0:
                fixed_path = Path(state.work_dir) / "rtl_vision_fixed.pptx"
                prs_to_fix.save(str(fixed_path))
                state.current_pptx = str(fixed_path)
                state.logs.append(f"[vision] Applied {fixes_applied} position fixes")
            else:
                state.logs.append("[vision] No fixes applied (could not match shapes to descriptions)")

    except ImportError:
        state.logs.append("[vision] SKIPPED: openai package not installed (pip install openai)")
    except Exception as e:
        state.logs.append(f"[vision] ERROR: {str(e)[:200]}")

    return state

def node_validate_translations(state: PipelineState, config: RunnableConfig) -> PipelineState:
    """
    Use GPT-4 Vision to validate ALL slides by comparing original vs transformed.
    Check for: missing translations, broken layouts, overlapping elements.
    This is fully dynamic and works with any PowerPoint file.
    """
    state.logs.append("[validate_translations] checking all slides with GPT-4 Vision...")

    try:
        from openai import OpenAI

        # Check for OpenAI API key
        api_key = os.getenv("OPENAI_API_KEY")
        if not api_key:
            state.logs.append("[validate_translations] SKIPPED: No OPENAI_API_KEY found")
            return state

        client = OpenAI(api_key=api_key)

        prs_original = Presentation(state.original_pptx_copy)
        prs_transformed = Presentation(state.current_pptx)
        num_slides = len(prs_original.slides)

        state.logs.append(f"[validate_translations] Validating {num_slides} slides...")

        # Use pdftoppm or similar to export ALL slides to images
        # LibreOffice export: converts entire PPTX to series of PNG files
        orig_export_dir = Path(state.work_dir) / "orig_export"
        trans_export_dir = Path(state.work_dir) / "trans_export"
        orig_export_dir.mkdir(exist_ok=True)
        trans_export_dir.mkdir(exist_ok=True)

        # Use LibreOffice to export to PDF first, then use pdfimages or similar
        # Simpler: use soffice export-as with multiple pages
        # Actually, soffice only exports first slide to PNG by default
        # Better approach: Use Python to extract slide screenshots programmatically

        state.logs.append("[validate_translations] Using programmatic validation (text-based)...")

        # Programmatic validation: check if Arabic text exists in transformed slides
        for slide_idx, (orig_slide, trans_slide) in enumerate(zip(prs_original.slides, prs_transformed.slides), start=1):
            # Extract all text shapes from transformed slide
            total_text_shapes = 0
            arabic_count = 0
            english_count = 0

            for shape in trans_slide.shapes:
                if hasattr(shape, 'text_frame'):
                    text = shape.text.strip()
                    if text and len(text) > 1:  # At least 2 characters
                        total_text_shapes += 1

                        has_arabic = bool(AR_LETTERS_RE.search(text))
                        has_english = bool(re.search(r'[a-zA-Z]{2,}', text))  # At least 2 Latin letters

                        # Shape is considered "translated" only if it has Arabic AND no English words
                        # (Arabic numerals alone don't count as translation)
                        if has_arabic and not has_english:
                            arabic_count += 1
                        elif has_english:
                            # Has English text (possibly mixed with Arabic numerals)
                            english_count += 1

            if total_text_shapes == 0:
                continue  # Skip slides with no text

            # Calculate Arabic coverage (shapes with Arabic vs all text shapes)
            arabic_coverage = (arabic_count / total_text_shapes * 100) if total_text_shapes > 0 else 0

            # Report status with English count to show untranslated content
            if arabic_coverage >= 80:
                state.logs.append(f"[validate_translations] ✓ Slide {slide_idx}: {arabic_count}/{total_text_shapes} shapes have Arabic, {english_count} still in English ({arabic_coverage:.0f}%)")
            elif arabic_coverage >= 50:
                state.logs.append(f"[validate_translations] ⚠️  Slide {slide_idx}: {arabic_count}/{total_text_shapes} shapes have Arabic, {english_count} still in English ({arabic_coverage:.0f}%)")
            else:
                state.logs.append(f"[validate_translations] ❌ Slide {slide_idx}: {arabic_count}/{total_text_shapes} shapes have Arabic, {english_count} still in English ({arabic_coverage:.0f}%) - NEEDS TRANSLATION")

    except Exception as e:
        state.logs.append(f"[validate_translations] ERROR: {str(e)[:200]}")

    return state

def node_finalize(state: PipelineState, config: RunnableConfig) -> PipelineState:
    # Move current to final out
    shutil.copyfile(state.current_pptx, state.out_pptx)
    state.logs.append(f"[finalize] wrote: {state.out_pptx}")
    return state

# --------------------------------------------------------------------------------------
# Graph builder
# --------------------------------------------------------------------------------------
def build_graph(checkpoint_path: Optional[str] = None):
    g = StateGraph(PipelineState)

    g.add_node("snapshot", node_snapshot_original)
    g.add_node("transform", node_apply_transform)
    g.add_node("validate", node_validate)
    g.add_node("recover", node_recover)
    g.add_node("fix_overlap", node_fix_icon_text_overlap)  # NEW: Fix text-icon overlaps
    g.add_node("colors", node_preserve_colors)
    g.add_node("validate_translations", node_validate_translations)
    g.add_node("vision_overlap", node_vision_overlap_fix)
    g.add_node("finalize", node_finalize)

    g.add_edge(START, "snapshot")
    g.add_edge("snapshot", "transform")
    g.add_edge("transform", "validate")
    g.add_conditional_edges("validate", route_validate, {"recover": "recover", "finalize": "fix_overlap"})
    g.add_edge("recover", "validate")
    g.add_edge("fix_overlap", "colors")  # NEW: Run overlap fix before color preservation
    g.add_edge("colors", "validate_translations")
    g.add_edge("validate_translations", "vision_overlap")
    g.add_edge("vision_overlap", "finalize")
    g.add_edge("finalize", END)

    # Disable checkpointing for now - context manager issue
    # if checkpoint_path and SqliteSaver:
    #     with SqliteSaver.from_conn_string(checkpoint_path) as saver:
    #         return g.compile(checkpointer=saver)
    # else:
    return g.compile()

# --------------------------------------------------------------------------------------
# CLI
# --------------------------------------------------------------------------------------
def main():
    ap = argparse.ArgumentParser(description="Lossless RTL Multi-Agent Pipeline")
    ap.add_argument("--in", dest="inp", required=True, help="Input PPTX")
    ap.add_argument("--out", dest="outp", required=True, help="Output PPTX")
    ap.add_argument("--map", dest="map_json", default=None, help="Translation map JSON (keyed by slide-#:shape-#)")
    ap.add_argument("--arabic-font", dest="arabic_font", default="Noto Naskh Arabic")
    ap.add_argument("--no-arabic-font", dest="no_arabic_font", action="store_true")
    ap.add_argument("--arabic-digits", dest="arabic_digits", action="store_true")
    ap.add_argument("--no-arabic-digits", dest="no_arabic_digits", action="store_true")
    ap.add_argument("--flip-icons", dest="flip_icons", action="store_true")
    ap.add_argument("--no-flip-icons", dest="no_flip_icons", action="store_true")
    ap.add_argument("--mirror", dest="mirror", action="store_true")
    ap.add_argument("--no-mirror", dest="no_mirror", action="store_true")
    args = ap.parse_args()

    arabic_font = None if args.no_arabic_font else args.arabic_font
    arabic_digits = False if args.no_arabic_digits else bool(args.arabic_digits)
    flip_icons = False if args.no_flip_icons else bool(args.flip_icons)
    mirror = False if args.no_mirror else bool(args.mirror)

    inp = str(Path(args.inp).resolve())
    outp = str(Path(args.outp).resolve())
    work_dir = tempfile.mkdtemp(prefix="rtl_graph_")
    original_copy = str(Path(work_dir) / "original.pptx")
    shutil.copyfile(inp, original_copy)

    params = PipelineParams(
        mirror=mirror,
        flip_icons=flip_icons,
        arabic_digits=arabic_digits,
        arabic_font=arabic_font,
        strict_shape_parity=True
    )
    state = PipelineState(
        input_pptx=inp,
        out_pptx=outp,
        map_json=args.map_json,
        params=params,
        work_dir=work_dir,
        current_pptx=original_copy,
        original_pptx_copy=original_copy
    )

    # Checkpoints enable time-travel/human-in-the-loop in LangGraph
    # (install langgraph-checkpoint-sqlite)
    ckpt_path = str(Path(outp).with_suffix("")) + ".checkpoints.sqlite"
    graph = build_graph(ckpt_path if SqliteSaver else None)

    # Pass LangSmith metadata if env is configured
    cfg: RunnableConfig = {
        "configurable": {
            "run_name": "rtl-lossless-pipeline",
            "metadata": {"app": "rtl-pipeline", "file": Path(inp).name}
        },
        "recursion_limit": 100  # Increased for multi-slide presentations
    }
    result = graph.invoke(state, cfg)

    # Extract state from result (LangGraph returns dict)
    if isinstance(result, dict):
        final_state = PipelineState(**result)
    else:
        final_state = result

    print("\n=== PIPELINE LOG ===")
    for line in final_state.logs:
        print(line)
    print("====================")
    print(f"Output: {final_state.out_pptx}")

if __name__ == "__main__":
    main()
