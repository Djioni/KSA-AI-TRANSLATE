#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Pixel-contrast Agent with Icon Features (no OCR):
- Renders slides to pixels (LibreOffice -> PDF -> PyMuPDF -> PNG).
- For each text shape, samples the rendered pixels inside the shape bbox.
- Uses Otsu threshold to separate foreground/background clusters.
- Computes WCAG contrast ratio; if below min threshold, recolors runs to a brand color that passes.
- Flips directional icons horizontally (arrows, chevrons, etc.) for RTL layout.
- Snaps icons to the right of nearest text (RTL convention).
- Writes an audit JSON with per-shape measurements and decisions.

Requirements:
  pip install python-pptx lxml pillow numpy pymupdf
  brew install --cask libreoffice   # for headless PPTX->PDF export

Usage:
  python pixel_contrast_agent.py \
    --in slides_AR.pptx \
    --out slides_AR_polished.pptx \
    --brand-dark "#0D2A47" \
    --brand-light "#FFFFFF" \
    --min-contrast 4.5 \
    --dpi 300 \
    --pad 6 \
    --flip-icons \
    --snap-icons \
    --audit-out audit_pixel.json

Notes:
- This runs AFTER your translation/RTL agent.
- Icon flipping is optional (use --flip-icons flag).
- Icon snapping is optional (use --snap-icons flag).
"""

from __future__ import annotations

import argparse
import json
import math
import os
import re
import shutil
import subprocess
import sys
import tempfile
from dataclasses import dataclass, asdict
from pathlib import Path
from typing import List, Optional, Tuple

import numpy as np
from PIL import Image
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.enum.lang import MSO_LANGUAGE_ID
from pptx.oxml.xmlchemy import OxmlElement

EMU_PER_INCH = 914400

# ------------------ Color & Contrast utils ------------------

def hex_to_rgb_tuple(s: str) -> Tuple[int, int, int]:
    s = s.strip().lstrip("#")
    if len(s) == 3:
        s = "".join([c*2 for c in s])
    return (int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))

def srgb_to_linear(c: float) -> float:
    # c in [0,1]
    return c/12.92 if c <= 0.04045 else ((c + 0.055)/1.055) ** 2.4

def rel_luminance(rgb: Tuple[int,int,int]) -> float:
    r, g, b = [x/255.0 for x in rgb]
    R = srgb_to_linear(r)
    G = srgb_to_linear(g)
    B = srgb_to_linear(b)
    return 0.2126*R + 0.7152*G + 0.0722*B

def contrast_ratio(fg: Tuple[int,int,int], bg: Tuple[int,int,int]) -> float:
    L1 = rel_luminance(fg)
    L2 = rel_luminance(bg)
    Lmax, Lmin = (L1, L2) if L1 >= L2 else (L2, L1)
    return (Lmax + 0.05) / (Lmin + 0.05)

# ------------------ Icon & Shape utils ------------------

# Name patterns for directional icons and logos
DIR_NAME_RE = re.compile(r"(arrow|chevron|caret|triangle-(?:right|left)|play|next|prev|bullet)", re.I)
LOGO_RE     = re.compile(r"(logo|brand|qrcode)", re.I)

def is_directional(shape_name: str) -> bool:
    """Check if shape name indicates a directional icon (arrow, chevron, etc.)"""
    return DIR_NAME_RE.search(shape_name or "") is not None

def is_logo_like(shape_name: str) -> bool:
    """Check if shape name indicates a logo/brand element that should not be flipped"""
    return LOGO_RE.search(shape_name or "") is not None

def flip_h(shape):
    """Flip shape horizontally by setting a:xfrm @flipH="1" in DrawingML"""
    try:
        # Find or create a:xfrm and set flipH="1"
        for xp in (".//a:xfrm", ".//p:spPr/a:xfrm", ".//p:grpSpPr/a:xfrm"):
            res = shape._element.xpath(xp, namespaces=shape._element.nsmap)
            if res:
                res[0].set("flipH", "1")
                return
        spPr = shape._element.xpath(".//p:spPr", namespaces=shape._element.nsmap)
        if spPr:
            xfrm = OxmlElement("a:xfrm")
            xfrm.set("flipH", "1")
            spPr[0].insert(0, xfrm)
    except Exception:
        pass

def bbox(shape) -> Tuple[int,int,int,int]:
    """Get bounding box of shape in EMU: (left, top, right, bottom)"""
    return int(shape.left), int(shape.top), int(shape.left + shape.width), int(shape.top + shape.height)

def y_overlap(a, b) -> int:
    """Calculate vertical overlap between two bboxes in EMU"""
    ay1, ay2 = a[1], a[3]
    by1, by2 = b[1], b[3]
    return max(0, min(ay2, by2) - max(ay1, by1))

# ------------------ Image rendering ------------------

def ensure_soffice() -> str:
    soffice = shutil.which("soffice")
    if not soffice:
        raise RuntimeError("LibreOffice `soffice` not found on PATH. Install via Homebrew: brew install --cask libreoffice")
    return soffice

def pptx_to_pdf(pptx_path: Path, out_dir: Path) -> Path:
    soffice = ensure_soffice()
    out_dir.mkdir(parents=True, exist_ok=True)
    cp = subprocess.run([soffice, "--headless", "--convert-to", "pdf", "--outdir", str(out_dir), str(pptx_path)],
                        capture_output=True, text=True)
    if cp.returncode != 0:
        raise RuntimeError(f"LibreOffice export failed:\nSTDOUT:\n{cp.stdout}\nSTDERR:\n{cp.stderr}")
    # find the produced pdf
    pdf_path = out_dir / (pptx_path.stem + ".pdf")
    if not pdf_path.exists():
        # fallback: first pdf in dir
        pdfs = list(out_dir.glob("*.pdf"))
        if not pdfs:
            raise RuntimeError("LibreOffice export produced no PDF.")
        pdf_path = pdfs[0]
    return pdf_path

def render_pdf_pages(pdf_path: Path, dpi: int) -> List[Image.Image]:
    doc = fitz.open(str(pdf_path))
    images: List[Image.Image] = []
    try:
        zoom = dpi / 72.0
        mat = fitz.Matrix(zoom, zoom)
        for i in range(doc.page_count):
            page = doc.load_page(i)
            pix = page.get_pixmap(matrix=mat, alpha=False)
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            images.append(img)
        return images
    finally:
        doc.close()

# ------------------ Geometry mapping ------------------

def emu_to_px(emu: int, dpi: int) -> int:
    inches = emu / EMU_PER_INCH
    return int(round(inches * dpi))

@dataclass
class SlideRenderInfo:
    slide_index: int  # 0-based
    img: Image.Image
    px_w: int
    px_h: int
    dpi: int

# ------------------ Otsu segmentation ------------------

def otsu_threshold(gray: np.ndarray) -> int:
    # gray: uint8 array
    hist = np.bincount(gray.ravel(), minlength=256).astype(np.float64)
    total = gray.size
    sum_all = (np.arange(256) * hist).sum()
    sum_b = 0.0
    w_b = 0.0
    max_var = -1.0
    threshold = 127
    for t in range(256):
        w_b += hist[t]
        if w_b == 0:
            continue
        w_f = total - w_b
        if w_f == 0:
            break
        sum_b += t * hist[t]
        m_b = sum_b / w_b
        m_f = (sum_all - sum_b) / w_f
        # between-class variance
        var_between = w_b * w_f * (m_b - m_f) ** 2
        if var_between > max_var:
            max_var = var_between
            threshold = t
    return int(threshold)

def estimate_fg_bg_from_region(region_rgb: np.ndarray) -> Tuple[Tuple[int,int,int], Tuple[int,int,int]]:
    """
    Returns (fg_rgb, bg_rgb) estimated from region pixels using Otsu.
    Typically text is darker; we pick the cluster with lower mean intensity as foreground.
    region_rgb: HxWx3 uint8
    """
    gray = (0.2126*region_rgb[...,0] + 0.7152*region_rgb[...,1] + 0.0722*region_rgb[...,2]).astype(np.uint8)
    t = otsu_threshold(gray)
    mask_dark = gray <= t
    # If mask is too small/large, flip strategy
    pct = mask_dark.mean()
    if pct < 0.02 or pct > 0.98:
        # not bimodal -> fallback: pick darkest 10% as fg
        thresh = np.percentile(gray, 10)
        mask_dark = gray <= thresh

    # Foreground = darker cluster
    fg_pixels = region_rgb[mask_dark]
    bg_pixels = region_rgb[~mask_dark]
    if fg_pixels.size == 0 or bg_pixels.size == 0:
        # degenerate; fallback to whole median as bg and black as fg
        median = np.median(region_rgb.reshape(-1,3), axis=0).astype(np.uint8)
        return (0,0,0), (int(median[0]), int(median[1]), int(median[2]))

    fg_median = np.median(fg_pixels, axis=0).astype(np.uint8)
    bg_median = np.median(bg_pixels, axis=0).astype(np.uint8)

    # ensure darker is fg
    if rel_luminance(tuple(fg_median)) > rel_luminance(tuple(bg_median)):
        fg_median, bg_median = bg_median, fg_median

    return (int(fg_median[0]), int(fg_median[1]), int(fg_median[2])), (int(bg_median[0]), int(bg_median[1]), int(bg_median[2]))

# ------------------ PPT helpers ------------------

def ensure_paragraph_rtl(tf):
    for p in tf.paragraphs:
        p.alignment = PP_ALIGN.RIGHT
        pPr = p._p.get_or_add_pPr()
        pPr.set("rtl", "1")
        pPr.set("algn", "r")
        for r in p.runs:
            try:
                r.font.language_id = MSO_LANGUAGE_ID.ARABIC_SAUDI_ARABIA
            except Exception:
                pass

def set_runs_color(tf, rgb: Tuple[int,int,int]):
    for p in tf.paragraphs:
        for r in p.runs:
            r.font.color.rgb = RGBColor(*rgb)

# ------------------ Audit ------------------

@dataclass
class AuditItem:
    slide: int
    shape_id: int
    name: str
    bbox_px: Tuple[int,int,int,int]
    measured_fg: Tuple[int,int,int]
    measured_bg: Tuple[int,int,int]
    ratio_before: float
    ratio_after: float
    fixed_contrast: bool
    applied_color: Optional[Tuple[int,int,int]]
    flipped_icon: bool = False
    snapped_icon: bool = False
    note: str = ""

# ------------------ Main logic ------------------

def process_pptx(in_path: Path, out_path: Path, brand_dark: Tuple[int,int,int], brand_light: Tuple[int,int,int],
                 min_contrast: float, dpi: int, pad_px: int, flip_icons: bool, snap_icons: bool,
                 icon_margin_emu: int, audit_out: Optional[Path]) -> None:
    prs = Presentation(str(in_path))
    slide_w_emu = int(prs.slide_width)
    slide_h_emu = int(prs.slide_height)

    with tempfile.TemporaryDirectory(prefix="px_contrast_") as td:
        td = Path(td)
        pdf_path = pptx_to_pdf(in_path, td)
        images = render_pdf_pages(pdf_path, dpi)
        assert len(images) == len(prs.slides), "Rendered page count differs from slide count."

        audits: List[AuditItem] = []

        for idx, slide in enumerate(prs.slides):
            img = images[idx]
            px_w, px_h = img.width, img.height
            # sanity check: expected pixels from EMU at given DPI
            exp_w = emu_to_px(slide_w_emu, dpi)
            exp_h = emu_to_px(slide_h_emu, dpi)
            scale_x = px_w / max(1, exp_w)
            scale_y = px_h / max(1, exp_h)

            # numpy copy once
            img_np = np.array(img)  # HxWx3

            # Collect text shapes and icon candidates for icon snapping
            text_shapes = [s for s in slide.shapes if getattr(s, "has_text_frame", False) and s.has_text_frame]
            icon_candidates = [s for s in slide.shapes
                             if s.shape_type in (MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.FREEFORM)]

            # Process text shapes for contrast fixes
            for shp in text_shapes:
                name = (getattr(shp, "name", "") or "").strip()

                # RTL safety (idempotent)
                ensure_paragraph_rtl(shp.text_frame)

                # EMU -> px bbox
                left_px = int(((int(shp.left)/EMU_PER_INCH) * dpi) * scale_x)
                top_px  = int(((int(shp.top) /EMU_PER_INCH) * dpi) * scale_y)
                w_px    = int(((int(shp.width)/EMU_PER_INCH) * dpi) * scale_x)
                h_px    = int(((int(shp.height)/EMU_PER_INCH)* dpi) * scale_y)

                x1 = max(0, left_px - pad_px)
                y1 = max(0, top_px  - pad_px)
                x2 = min(px_w, left_px + w_px + pad_px)
                y2 = min(px_h, top_px  + h_px + pad_px)

                region = img_np[y1:y2, x1:x2, :]
                if region.size == 0:
                    continue

                # Estimate fg/bg from pixels
                fg_rgb, bg_rgb = estimate_fg_bg_from_region(region)
                # If runs have explicit color, use the median of those as fg_before to measure "before"
                # Otherwise, use fg_rgb from pixels as proxy for before color.
                # We'll compute ratio_before using current run color if present.
                # Gather current fg colors:
                current_colors: List[Tuple[int,int,int]] = []
                for p in shp.text_frame.paragraphs:
                    for r in p.runs:
                        try:
                            c = r.font.color.rgb
                            if c is not None:
                                current_colors.append((c[0], c[1], c[2]))
                        except Exception:
                            pass
                if current_colors:
                    # median of current colors
                    cc = np.median(np.array(current_colors), axis=0).astype(np.uint8)
                    fg_before = (int(cc[0]), int(cc[1]), int(cc[2]))
                else:
                    fg_before = fg_rgb  # proxy

                ratio_before = contrast_ratio(fg_before, bg_rgb)

                applied = None
                ratio_after = ratio_before
                fixed = False
                note = ""

                if ratio_before < min_contrast:
                    # try brand_dark and brand_light; pick better
                    cr_dark  = contrast_ratio(brand_dark, bg_rgb)
                    cr_light = contrast_ratio(brand_light, bg_rgb)
                    best_rgb = brand_dark if cr_dark >= cr_light else brand_light
                    best_cr  = max(cr_dark, cr_light)

                    if best_cr > ratio_before:
                        set_runs_color(shp.text_frame, best_rgb)
                        applied = best_rgb
                        ratio_after = best_cr
                        fixed = True
                    else:
                        # As a fallback, force black/white whichever is best
                        bw_dark = (0,0,0); bw_light=(255,255,255)
                        cr_b = contrast_ratio(bw_dark, bg_rgb)
                        cr_w = contrast_ratio(bw_light, bg_rgb)
                        best = bw_dark if cr_b >= cr_w else bw_light
                        set_runs_color(shp.text_frame, best)
                        applied = best
                        ratio_after = max(cr_b, cr_w)
                        fixed = True
                        note = "used BW fallback"

                audits.append(AuditItem(
                    slide=idx+1,
                    shape_id=shp.shape_id,
                    name=name,
                    bbox_px=(x1,y1,x2,y2),
                    measured_fg=fg_before,
                    measured_bg=bg_rgb,
                    ratio_before=float(round(ratio_before,3)),
                    ratio_after=float(round(ratio_after,3)),
                    fixed_contrast=fixed,
                    applied_color=applied,
                    note=note
                ))

            # Process icons: flip directional icons
            if flip_icons:
                for icon in icon_candidates:
                    name = (getattr(icon, "name", "") or "").strip()
                    if not is_logo_like(name) and is_directional(name):
                        flip_h(icon)
                        audits.append(AuditItem(
                            slide=idx+1,
                            shape_id=icon.shape_id,
                            name=name,
                            bbox_px=(0,0,0,0),  # Not rendering icons
                            measured_fg=(0,0,0),
                            measured_bg=(0,0,0),
                            ratio_before=0.0,
                            ratio_after=0.0,
                            fixed_contrast=False,
                            applied_color=None,
                            flipped_icon=True,
                            note="directional icon"
                        ))

            # Snap icons to right of nearest text (RTL convention)
            if snap_icons and text_shapes:
                for icon in icon_candidates:
                    name = (getattr(icon, "name", "") or "").strip()
                    if is_logo_like(name):
                        continue

                    ib = bbox(icon)
                    # Find nearest text by vertical overlap
                    best = None
                    best_overlap = 0
                    for t in text_shapes:
                        tb = bbox(t)
                        ov = y_overlap(ib, tb)
                        if ov > best_overlap:
                            best_overlap = ov
                            best = t

                    if best and best_overlap > 0:
                        tb = bbox(best)
                        # Place icon to the RIGHT edge of the text shape (RTL start)
                        try:
                            new_left = tb[2] + icon_margin_emu
                            icon.left = new_left
                            audits.append(AuditItem(
                                slide=idx+1,
                                shape_id=icon.shape_id,
                                name=name,
                                bbox_px=(0,0,0,0),
                                measured_fg=(0,0,0),
                                measured_bg=(0,0,0),
                                ratio_before=0.0,
                                ratio_after=0.0,
                                fixed_contrast=False,
                                applied_color=None,
                                snapped_icon=True,
                                note=f"snapped to shape {best.shape_id}"
                            ))
                        except Exception as e:
                            pass

        prs.save(str(out_path))
        if audit_out:
            audit_out.write_text(json.dumps([asdict(a) for a in audits], ensure_ascii=False, indent=2), encoding="utf-8")

def main(argv=None):
    ap = argparse.ArgumentParser(description="Pixel-contrast agent with icon features for Arabic RTL slides")
    ap.add_argument("--in", dest="inp", required=True, help="Input PPTX (after RTL transform)")
    ap.add_argument("--out", dest="out", required=True, help="Output PPTX")
    ap.add_argument("--brand-dark", default="#0D2A47", help="Hex RGB for dark brand color (e.g., #0D2A47)")
    ap.add_argument("--brand-light", default="#FFFFFF", help="Hex RGB for light brand color")
    ap.add_argument("--min-contrast", type=float, default=4.5, help="Minimum WCAG contrast ratio to enforce")
    ap.add_argument("--dpi", type=int, default=300, help="Render DPI (240–300 recommended)")
    ap.add_argument("--pad", type=int, default=6, help="Padding (px) to sample around shape bbox")
    ap.add_argument("--flip-icons", action="store_true", help="Flip directional icons (arrows, chevrons, etc.)")
    ap.add_argument("--snap-icons", action="store_true", help="Snap icons to the RIGHT of nearest text")
    ap.add_argument("--icon-margin-emu", type=int, default=80000, help="Gap between text and snapped icon (~7pt)")
    ap.add_argument("--audit-out", default=None, help="Write audit JSON with before/after measurements")
    args = ap.parse_args(argv)

    in_path = Path(args.inp)
    out_path = Path(args.out)
    brand_dark = hex_to_rgb_tuple(args.brand_dark)
    brand_light = hex_to_rgb_tuple(args.brand_light)
    audit_out = Path(args.audit_out) if args.audit_out else None

    process_pptx(in_path, out_path, brand_dark, brand_light,
                 float(args.min_contrast), int(args.dpi), int(args.pad),
                 bool(args.flip_icons), bool(args.snap_icons), int(args.icon_margin_emu),
                 audit_out)
    print(f"✅ Wrote {out_path}")

if __name__ == "__main__":
    sys.exit(main())
