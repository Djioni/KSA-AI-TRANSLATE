#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Deterministic RTL transformer for PowerPoint (.pptx).

Features:
- Dump a translation map (JSON) with all text-bearing shapes and table cells.
- Apply translations (optional), enforce true RTL (a:pPr@rtl), right-align paragraphs,
  mirror X positions, flip directional icons/pictures, reverse table columns,
  enforce Arabic-capable font, optional Arabic-Indic digits, and simple contrast fix.
- Works without GPU or ML — pure Open XML + python-pptx.

CLI
  Dump map:
    python rtl_pptx_transformer.py dump-map input.pptx --out map.json
  Transform:
    python rtl_pptx_transformer.py transform input.pptx --map map.json --out output_AR.pptx \
      --flip-icons --arabic-font "Noto Naskh Arabic" --arabic-digits

Notes:
- Paragraph RTL must be set at paragraph level (a:pPr @rtl="1"), not only alignment.
- Horizontal flipping uses a:xfrm @flipH="1" in DrawingML.
"""

from __future__ import annotations
import argparse, json, sys
from pathlib import Path
from typing import Dict, Optional, Tuple

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.enum.lang import MSO_LANGUAGE_ID
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Pt
from pptx.dml.color import RGBColor

# ---------------- Utilities ----------------

def to_arabic_digits(s: str) -> str:
    trans = str.maketrans("0123456789", "٠١٢٣٤٥٦٧٨٩")
    return s.translate(trans)

def rel_key(slide_idx: int, shape_id: int) -> str:
    return f"slide-{slide_idx+1}:shape-{shape_id}"

def luminance(rgb: Tuple[int,int,int]) -> float:
    r, g, b = [x/255.0 for x in rgb]
    def lin(c): return c/12.92 if c <= 0.03928 else ((c+0.055)/1.055)**2.4
    R, G, B = lin(r), lin(g), lin(b)
    return 0.2126*R + 0.7152*G + 0.0722*B

def contrast_ratio(fg: Tuple[int,int,int], bg: Tuple[int,int,int]) -> float:
    L1, L2 = sorted([luminance(fg), luminance(bg)], reverse=True)
    return (L1 + 0.05) / (L2 + 0.05)

def pick_text_color(bg: Tuple[int,int,int]) -> Tuple[int,int,int]:
    black = (0,0,0); white=(255,255,255)
    # choose better contrast among black/white
    return black if contrast_ratio(black, bg) >= contrast_ratio(white, bg) else white

def shape_bg_rgb(shape) -> Tuple[int,int,int]:
    try:
        if shape.fill and shape.fill.type and shape.fill.fore_color and shape.fill.fore_color.rgb:
            rgb = shape.fill.fore_color.rgb
            return (rgb[0], rgb[1], rgb[2])
    except Exception:
        pass
    return (255,255,255)

def ensure_paragraph_rtl(p):
    # alignment
    p.alignment = PP_ALIGN.RIGHT
    # set underlying XML a:pPr rtl="1" and algn="r"
    pPr = p._p.get_or_add_pPr()
    pPr.set("rtl", "1")
    pPr.set("algn", "r")

def set_run_lang_and_font(run, family: Optional[str]):
    try:
        run.font.language_id = MSO_LANGUAGE_ID.ARABIC_SAUDI_ARABIA
    except Exception:
        pass
    if family:
        try:
            run.font.name = family
        except Exception:
            pass

def update_text_frame(tf, arabic_font: Optional[str], arabic_digits: bool):
    for p in tf.paragraphs:
        ensure_paragraph_rtl(p)
        for r in p.runs:
            if arabic_digits:
                r.text = to_arabic_digits(r.text)
            set_run_lang_and_font(r, arabic_font)

def set_text_color(shape, rgb: Tuple[int,int,int]):
    if not getattr(shape, "has_text_frame", False):
        return
    for p in shape.text_frame.paragraphs:
        for r in p.runs:
            r.font.color.rgb = RGBColor(*rgb)

def get_xfrm_element(sp_element):
    # try common locations; fall back to creating a:xfrm under spPr
    for xp in (".//a:xfrm", ".//p:spPr/a:xfrm", ".//p:grpSpPr/a:xfrm"):
        res = sp_element.xpath(xp, namespaces=sp_element.nsmap)
        if res:
            return res[0]
    spPr = sp_element.xpath(".//p:spPr", namespaces=sp_element.nsmap)
    if spPr:
        xfrm = OxmlElement("a:xfrm")
        spPr[0].insert(0, xfrm)
        return xfrm
    return None

def flip_h(shape):
    try:
        xfrm = get_xfrm_element(shape._element)
        if xfrm is not None:
            xfrm.set("flipH", "1")
    except Exception:
        pass

def is_directional(shape) -> bool:
    name = (getattr(shape, "name", "") or "").lower()
    return any(k in name for k in ("arrow", "chevron", "caret", "triangle-right", "play"))

def should_skip_flip(shape) -> bool:
    name = (getattr(shape, "name", "") or "").lower()
    return any(k in name for k in ("logo", "brand", "qrcode"))

def mirror_left(left: int, width: int, container_width: int) -> int:
    return int(container_width - (left + width))

def bbox(shape):
    return int(shape.left), int(shape.top), int(shape.left + shape.width), int(shape.top + shape.height)

def intersects(a, b) -> bool:
    ax1, ay1, ax2, ay2 = a; bx1, by1, bx2, by2 = b
    return not (ax2 <= bx1 or bx2 <= ax1 or ay2 <= by1 or by2 <= ay1)

def nudge_overlaps(shapes):
    # greedy downward nudging to reduce overlaps
    step = 12700  # ~1pt in EMU (approx); small increments
    ordered = sorted([s for s in shapes if hasattr(s, "left")], key=lambda s: (int(s.top), int(s.left)))
    for i in range(len(ordered)):
        for j in range(i):
            if intersects(bbox(ordered[i]), bbox(ordered[j])):
                ordered[i].top = ordered[j].top + ordered[j].height + step

# ---------------- Core logic ----------------

def dump_translation_map(in_pptx: Path, out_json: Optional[Path] = None) -> Dict[str, str]:
    prs = Presentation(str(in_pptx))
    mapping: Dict[str, str] = {}

    def handle_shape(s, slide_idx: int):
        st = s.shape_type
        key = rel_key(slide_idx, s.shape_id)

        # text-bearing shapes
        if getattr(s, "has_text_frame", False) and s.has_text_frame:
            mapping[key] = s.text or ""

        # tables (capture each cell)
        if st == MSO_SHAPE_TYPE.TABLE:
            tbl = s.table
            rows, cols = len(tbl.rows), len(tbl.columns)
            for r in range(rows):
                for c in range(cols):
                    cell_key = f"{key}:table:r{r}c{c}"
                    mapping[cell_key] = tbl.cell(r, c).text or ""

        # groups: recurse
        if st == MSO_SHAPE_TYPE.GROUP:
            for ch in s.shapes:
                handle_shape(ch, slide_idx)

    for s_idx, slide in enumerate(prs.slides):
        for shp in slide.shapes:
            handle_shape(shp, s_idx)

    if out_json:
        Path(out_json).write_text(json.dumps(mapping, ensure_ascii=False, indent=2), encoding="utf-8")
    return mapping

def apply_translations_to_shape(s, slide_idx: int, translations: Dict[str, str], arabic_digits: bool):
    key = rel_key(slide_idx, s.shape_id)

    if getattr(s, "has_text_frame", False) and s.has_text_frame:
        if key in translations:
            txt = translations[key]
            if arabic_digits:
                txt = to_arabic_digits(txt)

            # PRESERVE FORMATTING: Save original run properties before clearing
            original_colors = []
            original_fonts = []
            original_sizes = []

            for para in s.text_frame.paragraphs:
                for run in para.runs:
                    # Save color
                    try:
                        if run.font.color and run.font.color.rgb:
                            original_colors.append(run.font.color.rgb)
                        else:
                            original_colors.append(None)
                    except:
                        original_colors.append(None)

                    # Save font
                    try:
                        original_fonts.append(run.font.name)
                    except:
                        original_fonts.append(None)

                    # Save size
                    try:
                        original_sizes.append(run.font.size)
                    except:
                        original_sizes.append(None)

            # Now clear and add new text
            s.text_frame.clear()
            p = s.text_frame.paragraphs[0]
            r = p.add_run()
            r.text = txt

            # RESTORE FORMATTING: Apply first saved color/font/size
            if original_colors and original_colors[0]:
                r.font.color.rgb = original_colors[0]
            if original_fonts and original_fonts[0]:
                r.font.name = original_fonts[0]
            if original_sizes and original_sizes[0]:
                r.font.size = original_sizes[0]

    if s.shape_type == MSO_SHAPE_TYPE.TABLE:
        tbl = s.table
        rows, cols = len(tbl.rows), len(tbl.columns)
        for r in range(rows):
            for c in range(cols):
                cell_key = f"{key}:table:r{r}c{c}"
                if cell_key in translations:
                    txt = translations[cell_key]
                    if arabic_digits:
                        txt = to_arabic_digits(txt)
                    cell_tf = tbl.cell(r, c).text_frame
                    cell_tf.clear()
                    p = cell_tf.paragraphs[0]
                    rr = p.add_run()
                    rr.text = txt

def reverse_table_columns(tbl):
    rows, cols = len(tbl.rows), len(tbl.columns)
    for r in range(rows):
        for c in range(cols // 2):
            a = tbl.cell(r, c)
            b = tbl.cell(r, cols - 1 - c)
            a_text, b_text = a.text, b.text
            a.text, b.text = b_text, a_text

def enforce_table_rtl(tbl, arabic_font: Optional[str], arabic_digits: bool):
    rows, cols = len(tbl.rows), len(tbl.columns)
    for r in range(rows):
        for c in range(cols):
            update_text_frame(tbl.cell(r, c).text_frame, arabic_font, arabic_digits)

def process_shape(s, slide_idx: int, container_w: int, translations: Dict[str,str],
                  flip_icons: bool, arabic_font: Optional[str], arabic_digits: bool, fix_contrast: bool, mirror_positions: bool):
    st = s.shape_type

    # Apply translations first (so later transforms carry translated text)
    apply_translations_to_shape(s, slide_idx, translations, arabic_digits)

    # Tables: reverse columns + enforce RTL
    if st == MSO_SHAPE_TYPE.TABLE:
        tbl = s.table
        if mirror_positions:
            reverse_table_columns(tbl)
        enforce_table_rtl(tbl, arabic_font, arabic_digits)

    # Groups: recurse with group width
    if st == MSO_SHAPE_TYPE.GROUP:
        gw = int(s.width)
        for ch in s.shapes:
            process_shape(ch, slide_idx, gw, translations, flip_icons, arabic_font, arabic_digits, fix_contrast, mirror_positions)
        # Mirror group itself in parent container
        if mirror_positions:
            s.left = mirror_left(int(s.left), int(s.width), container_w)
        return

    # Mirror position inside container
    if mirror_positions and hasattr(s, "left") and hasattr(s, "width"):
        s.left = mirror_left(int(s.left), int(s.width), container_w)

    # Text frames: enforce RTL + font; digits already handled
    if getattr(s, "has_text_frame", False) and s.has_text_frame:
        update_text_frame(s.text_frame, arabic_font, arabic_digits)
        if fix_contrast:
            bg = shape_bg_rgb(s)
            set_text_color(s, pick_text_color(bg))

    # Flip arrows/pictures for RTL
    if flip_icons:
        if st in (MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.FREEFORM):
            if not should_skip_flip(s) and (is_directional(s) or st == MSO_SHAPE_TYPE.PICTURE):
                flip_h(s)

def transform(in_pptx: Path, out_pptx: Path, translations_path: Optional[Path],
              flip_icons: bool, arabic_font: Optional[str],
              arabic_digits: bool, fix_contrast: bool, mirror_positions: bool):
    translations: Dict[str,str] = {}
    if translations_path and translations_path.exists():
        translations = json.loads(Path(translations_path).read_text(encoding="utf-8"))

    prs = Presentation(str(in_pptx))
    slide_w = int(prs.slide_width)

    for s_idx, slide in enumerate(prs.slides):
        for shp in slide.shapes:
            process_shape(shp, s_idx, slide_w, translations, flip_icons, arabic_font, arabic_digits, fix_contrast, mirror_positions)
        # second pass: gentle overlap nudging (only if mirroring)
        if mirror_positions:
            nudge_overlaps(slide.shapes)

    prs.save(str(out_pptx))

# ---------------- CLI ----------------

def main(argv=None):
    ap = argparse.ArgumentParser(prog="rtl_pptx_transformer")
    sub = ap.add_subparsers(dest="cmd", required=True)

    ap_dump = sub.add_parser("dump-map", help="Dump translation map JSON")
    ap_dump.add_argument("pptx", help="Input PPTX")
    ap_dump.add_argument("--out", default=None, help="Output JSON (defaults to stdout)")

    ap_tr = sub.add_parser("transform", help="Transform PPTX to RTL Arabic variant")
    ap_tr.add_argument("pptx", help="Input PPTX")
    ap_tr.add_argument("--map", dest="map_json", default=None, help="Translations JSON")
    ap_tr.add_argument("--out", default=None, help="Output PPTX (default: *_AR.pptx)")
    ap_tr.add_argument("--flip-icons", action="store_true", help="Flip arrows/pictures horizontally")
    ap_tr.add_argument("--arabic-font", default="Noto Naskh Arabic", help="Arabic font family")
    ap_tr.add_argument("--arabic-digits", action="store_true", help="Convert 0-9 → Arabic-Indic")
    ap_tr.add_argument("--no-contrast-fix", action="store_true", help="Disable text contrast auto-fix")
    ap_tr.add_argument("--no-mirror", action="store_true", help="Disable position mirroring (text-only RTL)")

    args = ap.parse_args(argv)

    if args.cmd == "dump-map":
        mapping = dump_translation_map(Path(args.pptx), None)
        if args.out:
            Path(args.out).write_text(json.dumps(mapping, ensure_ascii=False, indent=2), encoding="utf-8")
        else:
            sys.stdout.write(json.dumps(mapping, ensure_ascii=False, indent=2))
        return 0

    if args.cmd == "transform":
        in_path = Path(args.pptx)
        out_path = Path(args.out or in_path.with_name(in_path.stem + "_AR.pptx"))
        transform(
            in_path,
            out_path,
            Path(args.map_json) if args.map_json else None,
            flip_icons=bool(args.flip_icons),
            arabic_font=args.arabic_font,
            arabic_digits=bool(args.arabic_digits),
            fix_contrast=(not args.no_contrast_fix),
            mirror_positions=(not args.no_mirror),
        )
        print(str(out_path))
        return 0

    return 1

if __name__ == "__main__":
    sys.exit(main())