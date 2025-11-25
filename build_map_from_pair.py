#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Build translation map from a pair of PPTX files (source EN + target AR).
This is useful for testing/validation when you have a gold Arabic deck.

Usage:
  python build_map_from_pair.py \
    --src "Template for Translation slide 2.pptx" \
    --dst "Template for Translation_AR slide 2.pptx" \
    --out "translations.json"
"""

import json
import argparse
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu

def emu(v):
    return int(v) if isinstance(v, Emu) else int(v)

def iter_shapes(slide, container_width):
    """Depth-first order; good enough if source & target keep similar order"""
    for shp in slide.shapes:
        yield shp, container_width
        if shp.shape_type == MSO_SHAPE_TYPE.GROUP:
            w = emu(shp.width)
            for sub in shp.shapes:
                yield sub, w

def collect_texts(prs):
    """Collect all text shapes per slide with slide-#:shape-# keys"""
    out = []
    sw = emu(prs.slide_width)
    for i, slide in enumerate(prs.slides, start=1):
        items = []
        for shp, cw in iter_shapes(slide, sw):
            if getattr(shp, "has_text_frame", False):
                items.append((f"slide-{i}:shape-{shp.shape_id}", (shp.text or "").strip()))
        out.append(items)
    return out

def main():
    ap = argparse.ArgumentParser(description="Build translation map from PPTX pair")
    ap.add_argument("--src", required=True, help="Source (EN) PPTX")
    ap.add_argument("--dst", required=True, help="Target (AR) PPTX")
    ap.add_argument("--out", required=True, help="Output translations.json")
    args = ap.parse_args()

    src = Presentation(args.src)
    dst = Presentation(args.dst)

    s = collect_texts(src)
    d = collect_texts(dst)
    mapping = {}

    for si, (s_items, d_items) in enumerate(zip(s, d), start=1):
        # Naive: pair by order among text shapes
        for (s_key, s_txt), (_, d_txt) in zip(s_items, d_items):
            if d_txt and d_txt.strip():
                mapping[s_key] = d_txt.strip()

    with open(args.out, "w", encoding="utf-8") as f:
        json.dump(mapping, f, ensure_ascii=False, indent=2)

    print(f"✅ Wrote {args.out} with {len(mapping)} translation entries")
    print(f"   Source: {args.src}")
    print(f"   Target: {args.dst}")

if __name__ == "__main__":
    main()
