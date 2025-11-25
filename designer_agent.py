#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Designer Agent (QA + Fix) for Arabic RTL slides.

Run this AFTER your RTL/translation transformer.

What it does (safe, deterministic):
  - Keeps brand colors by default, but FIXES low-contrast text (configurable).
  - Normalizes paragraph bidi/right-align (idempotent).
  - Snaps directional icons next to the nearest text on the RIGHT (RTL convention).
  - Optionally flips only directional icons (arrows/chevrons/play/etc.).
  - Gently nudges shapes to reduce overlaps.

Dependencies:
  pip install python-pptx lxml

Usage:
  python designer_agent.py \
    --in slides_AR.pptx \
    --out slides_AR_polished.pptx \
    --brand-dark "#0D2A47" \
    --brand-light "#FFFFFF" \
    --min-contrast 4.5 \
    --flip-directional-icons \
    --snap-icons

Notes
- We preserve existing explicit run colors unless contrast is *really* poor.
- Icon snapping is based on nearest text shape by Y-overlap; margin is configurable.
"""

from __future__ import annotations

import argparse
import json
import sys
from dataclasses import dataclass, asdict
from pathlib import Path
from typing import List, Optional, Tuple
import re

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.enum.lang import MSO_LANGUAGE_ID
from pptx.oxml.xmlchemy import OxmlElement
from pptx.dml.color import RGBColor

# -------- Utils

def hex_to_rgb_tuple(s: str) -> Tuple[int, int, int]:
    s = s.strip().lstrip("#")
    if len(s) == 3:
        s = "".join([c*2 for c in s])
    return (int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))

def luminance(rgb: Tuple[int,int,int]) -> float:
    r, g, b = [x/255.0 for x in rgb]
    def lin(c): return c/12.92 if c <= 0.03928 else ((c+0.055)/1.055)**2.4
    R, G, B = lin(r), lin(g), lin(b)
    return 0.2126*R + 0.7152*G + 0.0722*B

def contrast_ratio(fg: Tuple[int,int,int], bg: Tuple[int,int,int]) -> float:
    L1, L2 = sorted([luminance(fg), luminance(bg)], reverse=True)
    return (L1 + 0.05) / (L2 + 0.05)

def get_shape_bg_rgb(shape) -> Tuple[int,int,int]:
    # Try the shape's fill; if none, assume white (canvas)
    try:
        if shape.fill and shape.fill.type and shape.fill.fore_color and shape.fill.fore_color.rgb:
            c = shape.fill.fore_color.rgb
            return (c[0], c[1], c[2])
    except Exception:
        pass
    return (255, 255, 255)

def ensure_paragraph_rtl(p):
    p.alignment = PP_ALIGN.RIGHT
    pPr = p._p.get_or_add_pPr()
    pPr.set("rtl", "1")
    pPr.set("algn", "r")

def ensure_textframe_rtl(tf):
    for p in tf.paragraphs:
        ensure_paragraph_rtl(p)
        for r in p.runs:
            try:
                r.font.language_id = MSO_LANGUAGE_ID.ARABIC_SAUDI_ARABIA
            except Exception:
                pass

def set_run_rgb(run, rgb: Tuple[int,int,int]):
    run.font.color.rgb = RGBColor(*rgb)

def get_run_rgb(run) -> Optional[Tuple[int,int,int]]:
    try:
        if run.font.color and run.font.color.rgb:
            rgb = run.font.color.rgb
            return (rgb[0], rgb[1], rgb[2])
    except Exception:
        pass
    return None

def bbox(shape) -> Tuple[int,int,int,int]:
    return int(shape.left), int(shape.top), int(shape.left + shape.width), int(shape.top + shape.height)

def center_xy(shape) -> Tuple[int,int]:
    return int(shape.left) + int(shape.width)//2, int(shape.top) + int(shape.height)//2

def intersects(a, b) -> bool:
    ax1, ay1, ax2, ay2 = a; bx1, by1, bx2, by2 = b
    return not (ax2 <= bx1 or bx2 <= ax1 or ay2 <= by1 or by2 <= ay1)

def y_overlap(a, b) -> int:
    # amount of vertical overlap in EMU
    ay1, ay2 = a[1], a[3]
    by1, by2 = b[1], b[3]
    return max(0, min(ay2, by2) - max(ay1, by1))

# Name patterns
DIR_NAME_RE = re.compile(r"(arrow|chevron|caret|triangle-(?:right|left)|play|next|prev|bullet)", re.I)
LOGO_RE     = re.compile(r"(logo|brand|qrcode)", re.I)

def is_directional(shape_name: str) -> bool:
    return DIR_NAME_RE.search(shape_name or "") is not None

def is_logo_like(shape_name: str) -> bool:
    return LOGO_RE.search(shape_name or "") is not None

def flip_h(shape):
    try:
        # Find or create a:xfrm and set flipH="1"
        for xp in (".//a:xfrm", ".//p:spPr/a:xfrm", ".//p:grpSpPr/a:xfrm"):
            res = shape._element.xpath(xp, namespaces=shape._element.nsmap)
            if res:
                res[0].set("flipH", "1")
                return
        spPr = shape._element.xpath(".//p:spPr", namespaces=shape._element.nsmap)
        if spPr:
            xfrm = OxmlElement("a:xfrm")
            xfrm.set("flipH", "1")
            spPr[0].insert(0, xfrm)
    except Exception:
        pass

def nudge_overlaps(shapes):
    # Gentle downward nudging to reduce collisions
    step = 12700  # ~1pt EMU
    ordered = sorted([s for s in shapes if hasattr(s, "left")], key=lambda s: (int(s.top), int(s.left)))
    for i in range(len(ordered)):
        for j in range(i):
            if intersects(bbox(ordered[i]), bbox(ordered[j])):
                ordered[i].top = ordered[j].top + ordered[j].height + step

# -------- OCR Validation

def validate_with_ocr(pptx_path: Path, dpi: int = 150) -> dict:
    """
    Render slides to images and OCR them to verify text is readable.
    Returns: {"ok": bool, "slides": [{"slide": int, "readable": bool, "ocr_confidence": float, "sample_text": str}]}

    Requires: brew install tesseract
              pip install pytesseract pymupdf pillow
    """
    try:
        import fitz  # PyMuPDF
        from PIL import Image
        import pytesseract
        import subprocess
        import shutil
        import tempfile
    except ImportError as e:
        return {"ok": False, "error": f"Missing OCR dependencies: {e}"}

    # Convert PPTX to PDF using LibreOffice
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if not soffice:
        return {"ok": False, "error": "LibreOffice not found (needed for rendering). Install with: brew install --cask libreoffice"}

    with tempfile.TemporaryDirectory(prefix="ocr_validate_") as td:
        temp_dir = Path(td)
        pdf_path = temp_dir / "slides.pdf"

        # Convert PPTX to PDF
        result = subprocess.run(
            [soffice, "--headless", "--convert-to", "pdf", "--outdir", str(temp_dir), str(pptx_path)],
            capture_output=True,
            text=True,
            timeout=60
        )

        if result.returncode != 0:
            return {"ok": False, "error": f"PDF conversion failed: {result.stderr}"}

        # Find the generated PDF (LibreOffice names it based on input filename)
        pdf_files = list(temp_dir.glob("*.pdf"))
        if not pdf_files:
            return {"ok": False, "error": "No PDF generated"}

        pdf_path = pdf_files[0]

        # Open PDF and OCR each page
        doc = fitz.open(str(pdf_path))
        results = []

        for page_num in range(doc.page_count):
            page = doc.load_page(page_num)
            zoom = dpi / 72.0
            mat = fitz.Matrix(zoom, zoom)
            pix = page.get_pixmap(matrix=mat, alpha=False)

            # Save to PNG
            img_path = temp_dir / f"slide_{page_num + 1}.png"
            pix.save(str(img_path))

            # OCR with Tesseract (Arabic + English)
            try:
                img = Image.open(str(img_path))
                ocr_data = pytesseract.image_to_data(img, lang="ara+eng", output_type=pytesseract.Output.DICT)

                # Calculate average confidence for detected text
                confidences = [int(conf) for conf in ocr_data['conf'] if conf != '-1']
                avg_confidence = sum(confidences) / len(confidences) if confidences else 0

                # Get sample text
                text = pytesseract.image_to_string(img, lang="ara+eng").strip()
                sample = text[:200] if text else ""

                # Consider readable if: confidence > 60% AND some text detected
                readable = avg_confidence > 60 and len(text) > 20

                results.append({
                    "slide": page_num + 1,
                    "readable": readable,
                    "ocr_confidence": round(avg_confidence, 2),
                    "text_length": len(text),
                    "sample_text": sample
                })
            except Exception as e:
                results.append({
                    "slide": page_num + 1,
                    "readable": False,
                    "ocr_confidence": 0,
                    "error": str(e)
                })

        doc.close()

        # Overall success if all slides are readable
        all_readable = all(r.get("readable", False) for r in results)

        return {
            "ok": True,
            "all_readable": all_readable,
            "slides": results
        }

    return {"ok": False, "error": "Unknown error"}

# -------- Audit data

@dataclass
class FixLog:
    slide: int
    shape_id: int
    name: str
    fixed_contrast: bool
    snapped_icon: bool
    flipped_icon: bool
    rtl_enforced: bool
    notes: str = ""

# -------- Core fixer

def fix_slide(slide, slide_index: int, slide_w: int, slide_h: int,
              brand_dark: Tuple[int,int,int], brand_light: Tuple[int,int,int],
              min_contrast: float,
              flip_directional_icons: bool,
              snap_icons: bool,
              icon_margin_emu: int,
              audit: List[FixLog]):

    shapes = list(slide.shapes)

    # Precompute text shapes & icon candidates
    text_shapes = [s for s in shapes if getattr(s, "has_text_frame", False) and s.has_text_frame]
    icon_candidates = [s for s in shapes
                       if s.shape_type in (MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.FREEFORM)]

    # 1) Text: enforce RTL, then check contrast
    for s in text_shapes:
        name = (getattr(s, "name", "") or "").strip()
        log = FixLog(slide=slide_index, shape_id=s.shape_id, name=name,
                     fixed_contrast=False, snapped_icon=False, flipped_icon=False, rtl_enforced=False)

        # Enforce RTL paragraphs
        try:
            ensure_textframe_rtl(s.text_frame)
            log.rtl_enforced = True
        except Exception as e:
            log.notes += f"rtl_err:{e}; "

        # Contrast fix (conservative): only when EVERY run is too low-contrast vs background
        try:
            bg = get_shape_bg_rgb(s)
            all_low = True
            for p in s.text_frame.paragraphs:
                for r in p.runs:
                    fg = get_run_rgb(r) or brand_dark  # assume dark if not explicit
                    cr = contrast_ratio(fg, bg)
                    if cr >= min_contrast:
                        all_low = False
                        break
                if not all_low:
                    break

            if all_low:
                # If background is light, use brand_dark; if dark, use brand_light
                use = brand_dark if contrast_ratio(brand_dark, bg) >= contrast_ratio(brand_light, bg) else brand_light
                for p in s.text_frame.paragraphs:
                    for r in p.runs:
                        set_run_rgb(r, use)
                log.fixed_contrast = True
        except Exception as e:
            log.notes += f"contrast_err:{e}; "

        audit.append(log)

    # 2) Icons: optional horizontal flip (directional only, not logos)
    for s in icon_candidates:
        name = (getattr(s, "name", "") or "").strip()
        if flip_directional_icons and not is_logo_like(name) and is_directional(name):
            flip_h(s)
            audit.append(FixLog(slide=slide_index, shape_id=s.shape_id, name=name,
                                fixed_contrast=False, snapped_icon=False, flipped_icon=True, rtl_enforced=False))

    # 3) Icon snapping: put icons to the RIGHT of the nearest text (RTL "leading" side)
    if snap_icons and text_shapes:
        for icon in icon_candidates:
            name = (getattr(icon, "name", "") or "").strip()
            if is_logo_like(name):
                continue
            ib = bbox(icon)
            # Find nearest text by vertical overlap
            best = None
            best_overlap = 0
            for t in text_shapes:
                tb = bbox(t)
                ov = y_overlap(ib, tb)
                if ov > best_overlap:
                    best_overlap = ov
                    best = t
            if best and best_overlap > 0:
                tb = bbox(best)
                # Place icon to the RIGHT edge of the text shape (RTL start)
                try:
                    new_left = tb[2] + icon_margin_emu
                    icon.left = new_left
                    audit.append(FixLog(slide=slide_index, shape_id=icon.shape_id, name=name,
                                        fixed_contrast=False, snapped_icon=True, flipped_icon=False, rtl_enforced=False))
                except Exception:
                    pass

    # 4) Gentle overlap reduction
    # DISABLED - This moves shapes vertically and breaks the layout after RTL mirroring
    # nudge_overlaps(slide.shapes)

# -------- Main

def main(argv=None):
    ap = argparse.ArgumentParser(description="Designer Agent (QA + Fix) for RTL Arabic slides")
    ap.add_argument("--in", dest="inp", required=True, help="Input PPTX (from RTL transformer)")
    ap.add_argument("--out", dest="out", required=True, help="Output PPTX")
    ap.add_argument("--brand-dark", default="#0D2A47", help="Dark brand text color (hex)")
    ap.add_argument("--brand-light", default="#FFFFFF", help="Light brand text color (hex)")
    ap.add_argument("--min-contrast", type=float, default=4.5, help="Min contrast ratio for text vs background")
    ap.add_argument("--flip-directional-icons", action="store_true", help="Flip arrows/chevrons/play etc.")
    ap.add_argument("--snap-icons", action="store_true", help="Snap icons to the RIGHT of nearest text")
    ap.add_argument("--icon-margin-emu", type=int, default=80000, help="Gap between text and snapped icon (~7pt)")
    ap.add_argument("--audit-out", default=None, help="Write JSON of fixes applied (optional)")
    ap.add_argument("--ocr-validate", action="store_true", help="Run OCR validation after fixes")
    ap.add_argument("--ocr-report", default=None, help="Write OCR validation report JSON (optional)")
    args = ap.parse_args(argv)

    prs = Presentation(str(args.inp))
    slide_w = int(prs.slide_width)
    slide_h = int(prs.slide_height)

    brand_dark = hex_to_rgb_tuple(args.brand_dark)
    brand_light = hex_to_rgb_tuple(args.brand_light)

    audit: List[FixLog] = []

    for idx, slide in enumerate(prs.slides, start=1):
        fix_slide(
            slide, idx, slide_w, slide_h,
            brand_dark=brand_dark, brand_light=brand_light,
            min_contrast=float(args.min_contrast),
            flip_directional_icons=bool(args.flip_directional_icons),
            snap_icons=bool(args.snap_icons),
            icon_margin_emu=int(args.icon_margin_emu),
            audit=audit
        )

    prs.save(str(args.out))

    if args.audit_out:
        Path(args.audit_out).write_text(json.dumps([asdict(a) for a in audit], ensure_ascii=False, indent=2), encoding="utf-8")

    print(f"✅ Wrote {args.out}")

    # Optional OCR validation
    if args.ocr_validate or args.ocr_report:
        print("🔍 Running OCR validation...")
        ocr_result = validate_with_ocr(Path(args.out))

        if args.ocr_report:
            Path(args.ocr_report).write_text(json.dumps(ocr_result, ensure_ascii=False, indent=2), encoding="utf-8")
            print(f"📄 OCR report: {args.ocr_report}")

        if ocr_result.get("ok"):
            if ocr_result.get("all_readable"):
                print("✅ OCR validation passed - all slides readable")
            else:
                print("⚠️  OCR validation warning - some slides may have visibility issues")
                for slide in ocr_result.get("slides", []):
                    if not slide.get("readable"):
                        print(f"  Slide {slide['slide']}: confidence={slide.get('ocr_confidence', 0)}%")
        else:
            print(f"⚠️  OCR validation failed: {ocr_result.get('error', 'Unknown error')}")

    return 0

if __name__ == "__main__":
    sys.exit(main())
